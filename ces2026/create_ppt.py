#!/usr/bin/env python3
"""CES 2026 현장 리포트 — 비주얼 중심 1시간 발표용 PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

BASE = os.path.dirname(os.path.abspath(__file__))
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG = RGBColor(0x0B, 0x0B, 0x0F)
SURFACE = RGBColor(0x16, 0x16, 0x1E)
TEXT = RGBColor(0xE8, 0xE6, 0xF0)
MUTED = RGBColor(0x9C, 0x9A, 0xA8)
AI = RGBColor(0x6C, 0x63, 0xFF)
GREEN = RGBColor(0x10, 0xB9, 0x81)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x1E, 0x2E)
OVERLAY = RGBColor(0x0B, 0x0B, 0x0F)

def img(name):
    return os.path.join(BASE, name)

def set_bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def txt(slide, l, t, w, h, text, size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Arial'
    p.alignment = align
    return box

def multi_txt(slide, l, t, w, h, lines_data):
    """lines_data = [(text, size, color, bold), ...]"""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold) in enumerate(lines_data):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = 'Arial'
        p.space_after = Pt(4)
    return box

def box(slide, l, t, w, h, color=SURFACE):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s

def overlay(slide, l, t, w, h, alpha=0.7):
    """Dark overlay for text readability on photos"""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = OVERLAY
    s.line.fill.background()
    s.shadow.inherit = False
    return s

def photo_slide(image_file, title, subtitle="", title_size=40):
    """Full-bleed photo placeholder with text overlay"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    # Photo placeholder
    ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(4.0))
    ph.fill.solid()
    ph.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x30)
    ph.line.color.rgb = RGBColor(0x3A, 0x3A, 0x4A)
    ph.line.width = Pt(1)
    txt(slide, 0.5, 1.2, 12, 0.5, "📷 사진 삽입 위치", 16, AI, True, PP_ALIGN.CENTER)
    txt(slide, 0.5, 1.8, 12, 1.5, PHOTO_DESC.get(image_file, image_file), 14, MUTED, False, PP_ALIGN.CENTER)
    txt(slide, 0.8, 4.5, 11, 1.5, title, title_size, TEXT, True, PP_ALIGN.LEFT)
    if subtitle:
        txt(slide, 0.8, 5.8, 11, 1, subtitle, 20, MUTED, False, PP_ALIGN.LEFT)
    return slide

def photo_left(image_file, title, bullets, note=""):
    """Photo placeholder on left, content on right"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(6.5), prs.slide_height)
    ph.fill.solid()
    ph.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x30)
    ph.line.color.rgb = RGBColor(0x3A, 0x3A, 0x4A)
    ph.line.width = Pt(1)
    txt(slide, 0.3, 2.5, 5.9, 0.5, "📷 사진 삽입 위치", 14, AI, True, PP_ALIGN.CENTER)
    txt(slide, 0.3, 3.1, 5.9, 1.5, PHOTO_DESC.get(image_file, image_file), 12, MUTED, False, PP_ALIGN.CENTER)
    txt(slide, 7, 0.5, 5.8, 0.8, title, 26, TEXT, True)
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(1.2), Inches(1.2), Pt(3))
    s.fill.solid(); s.fill.fore_color.rgb = AI; s.line.fill.background()
    y = 1.6
    for b in bullets:
        if b == "":
            y += 0.2
            continue
        txt(slide, 7.2, y, 5.5, 0.5, b, 16, MUTED)
        y += 0.42
    if note:
        txt(slide, 7, 6.3, 5.8, 0.5, note, 14, AI, True)
    return slide

def photo_right(image_file, title, bullets, note=""):
    """Content on left, photo placeholder on right"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.833), 0, Inches(6.5), prs.slide_height)
    ph.fill.solid()
    ph.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x30)
    ph.line.color.rgb = RGBColor(0x3A, 0x3A, 0x4A)
    ph.line.width = Pt(1)
    txt(slide, 7.1, 2.5, 5.9, 0.5, "📷 사진 삽입 위치", 14, AI, True, PP_ALIGN.CENTER)
    txt(slide, 7.1, 3.1, 5.9, 1.5, PHOTO_DESC.get(image_file, image_file), 12, MUTED, False, PP_ALIGN.CENTER)
    txt(slide, 0.8, 0.5, 5.5, 0.8, title, 26, TEXT, True)
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(1.2), Pt(3))
    s.fill.solid(); s.fill.fore_color.rgb = AI; s.line.fill.background()
    y = 1.6
    for b in bullets:
        if b == "":
            y += 0.2
            continue
        txt(slide, 1.0, y, 5.3, 0.5, b, 16, MUTED)
        y += 0.42
    if note:
        txt(slide, 0.8, 6.3, 5.5, 0.5, note, 14, AI, True)
    return slide

def stat_card(slide, l, t, w, h, value, label):
    box(slide, l, t, w, h, SURFACE)
    txt(slide, l+0.1, t+0.3, w-0.2, 0.8, value, 32, AI, True, PP_ALIGN.CENTER)
    txt(slide, l+0.1, t+1.2, w-0.2, 0.6, label, 14, MUTED, False, PP_ALIGN.CENTER)

def section_div(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    txt(slide, 1, 2.5, 11, 1, title, 40, AI, True, PP_ALIGN.CENTER)
    if subtitle:
        txt(slide, 1.5, 3.8, 10, 1, subtitle, 22, MUTED, False, PP_ALIGN.CENTER)
    return slide

def headline_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    txt(slide, 1, 2.2, 11, 2, title, 36, TEXT, True, PP_ALIGN.CENTER)
    if subtitle:
        txt(slide, 1.5, 4.5, 10, 1.5, subtitle, 20, MUTED, False, PP_ALIGN.CENTER)
    return slide

def key_msg_slide(number, title, body):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    txt(slide, 1, 1.5, 2, 2, f"0{number}", 80, AI, True, PP_ALIGN.LEFT)
    txt(slide, 3.5, 1.8, 9, 1.2, title, 30, TEXT, True)
    txt(slide, 3.5, 3.2, 9, 3, body, 20, MUTED)
    return slide



# ============================================================
# SLIDES
# ============================================================

P = {
    "cover": "KakaoTalk_Photo_2026-03-10-00-03-32 014.jpeg",
    "hall": "KakaoTalk_Photo_2026-03-10-00-01-57 001.jpeg",
    "floor1": "KakaoTalk_Photo_2026-03-10-00-01-57 003.jpeg",
    "floor2": "KakaoTalk_Photo_2026-03-10-00-01-57 005.jpeg",
    "booth": "KakaoTalk_Photo_2026-03-10-00-01-58 007.jpeg",
    "zoox_ext": "KakaoTalk_Photo_2026-03-10-00-03-30 001.jpeg",
    "zoox_side": "KakaoTalk_Photo_2026-03-10-00-03-31 002.jpeg",
    "zoox_in1": "KakaoTalk_Photo_2026-03-10-00-03-31 005.jpeg",
    "zoox_in2": "KakaoTalk_Photo_2026-03-10-00-03-31 007.jpeg",
    "zoox_disp": "KakaoTalk_Photo_2026-03-10-00-03-31 009.jpeg",
    "zoox_ride": "KakaoTalk_Photo_2026-03-10-00-03-32 010.jpeg",
    "display": "KakaoTalk_Photo_2026-03-10-00-03-32 012.jpeg",
    "waymo1": "KakaoTalk_Photo_2026-03-10-00-03-32 017.jpeg",
    "waymo2": "KakaoTalk_Photo_2026-03-10-00-03-33 018.jpeg",
    "vegas": "KakaoTalk_Photo_2026-03-10-00-03-33 021.jpeg",
    "street1": "KakaoTalk_Photo_2026-03-10-00-03-33 023.jpeg",
    "street2": "KakaoTalk_Photo_2026-03-10-00-03-33 026.jpeg",
    "night": "KakaoTalk_Photo_2026-03-10-00-03-34 027.jpeg",
    "pano": "KakaoTalk_Photo_2026-03-10-00-03-34 029.jpeg",
}

# Photo descriptions for placeholders
PHOTO_DESC = {
    P["cover"]: "CES 2026 전시장 전경 사진\n(컨벤션 센터 입구 또는 메인 홀 전경)",
    P["hall"]: "CES 2026 웨스트 홀 입구 사진\n(참가자들이 입장하는 모습)",
    P["floor1"]: "CES 2026 쇼 플로어 사진\n(NVIDIA 또는 주요 부스가 보이는 전시장)",
    P["floor2"]: "CES 2026 쇼 플로어 전경\n(넓은 전시장 모습, 다양한 부스)",
    P["booth"]: "CES 2026 전시 부스 사진\n(로봇 또는 자동차 관련 부스)",
    P["zoox_ext"]: "Zoox 로보택시 외관 사진\n(핸들 없는 전용 설계 차량 전체 모습)",
    P["zoox_side"]: "Zoox 로보택시 측면 사진\n(양방향 좌석 구조가 보이는 각도)",
    P["zoox_in1"]: "Zoox 차량 내부 사진\n(양방향 좌석, 디스플레이, LiDAR 센서)",
    P["zoox_in2"]: "Zoox 내부 좌석 사진\n(마주보는 좌석 배치)",
    P["zoox_disp"]: "Zoox 내부 디스플레이 사진\n(온도/음악/경로 조작 화면)",
    P["zoox_ride"]: "Zoox 탑승 시점 사진\n(실제 탑승 중 촬영)",
    P["display"]: "CES 2026 디스플레이 전시 사진\n(Samsung/LG 등 디스플레이 부스)",
    P["waymo1"]: "Waymo 로보택시 사진\n(Jaguar I-PACE 기반 자율주행 차량)",
    P["waymo2"]: "Waymo 센서 어레이 사진\n(차량 상단 센서 장비 클로즈업)",
    P["vegas"]: "Vegas Loop 터널 사진\n(지하 터널 내부 또는 Tesla 차량)",
    P["street1"]: "라스베이거스 현장 사진\n(CES 기간 거리 풍경)",
    P["street2"]: "CES 2026 거리 풍경\n(라스베이거스 스트립 또는 컨벤션 주변)",
    P["night"]: "라스베이거스 야경 사진\n(CES 기간 밤 풍경)",
    P["pano"]: "CES 2026 전경 사진\n(파노라마 또는 넓은 각도)",
}



# --- 1. COVER ---
photo_slide(P["cover"], "CES 2026 참관 보고서",
    "Physical AI와 Software Defined Vehicle — AI가 화면을 벗어나 세상으로 나오다\n우지환 · 2026년 1월")

# --- 2. AGENDA ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
txt(slide, 0.8, 0.4, 5, 0.8, "Agenda", 28, TEXT, True)
items = [
    ("01", "들어가며 — 화면 밖으로 나온 지능"),
    ("02", "Physical AI — 지능이 몸을 얻다"),
    ("03", "SDV — 자동차의 OS 전환"),
    ("04", "현장 체험 — Zoox · Waymo · Vegas Loop"),
    ("05", "주목 기업 50"),
    ("06", "다섯 가지 메시지"),
]
y = 1.5
for num, label in items:
    txt(slide, 1.0, y, 1, 0.5, num, 24, AI, True)
    txt(slide, 2.0, y, 5, 0.5, label, 20, MUTED)
    y += 0.7
ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(0.5), Inches(5.8), Inches(6.5))
ph.fill.solid(); ph.fill.fore_color.rgb = RGBColor(0x22,0x22,0x30); ph.line.color.rgb = RGBColor(0x3A,0x3A,0x4A); ph.line.width = Pt(1)
txt(slide, 7.3, 3, 5.2, 1, "📷 CES 2026 전경 사진\n(라스베이거스 야경 또는 컨벤션 센터 전경)", 12, MUTED, False, PP_ALIGN.CENTER)

# --- 3. 들어가며 ---
photo_slide(P["hall"], "CES 2026은 달랐다",
    "로봇이 걷고, 차가 스스로 생각하고, AI가 물리 세계로 나왔다")

# --- 4. CES 규모 ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
txt(slide, 0.8, 0.4, 11, 0.8, "CES 2026 한눈에 보기", 28, TEXT, True)
stat_card(slide, 0.8, 1.5, 2.8, 2, "148K+", "참가자")
stat_card(slide, 4.0, 1.5, 2.8, 2, "4,500+", "전시 기업")
stat_card(slide, 7.2, 1.5, 2.8, 2, "250만ft²", "전시 면적")
stat_card(slide, 10.4, 1.5, 2.8, 2, "13개", "베뉴")
ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4), Inches(11.7), Inches(3.2))
ph.fill.solid(); ph.fill.fore_color.rgb = RGBColor(0x22,0x22,0x30); ph.line.color.rgb = RGBColor(0x3A,0x3A,0x4A); ph.line.width = Pt(1)
txt(slide, 1, 5, 11, 0.8, "📷 CES 2026 쇼 플로어 전경 사진 (넓은 전시장 모습)", 12, MUTED, False, PP_ALIGN.CENTER)

# --- 5. Jensen Huang 선언 ---
photo_right(P["floor1"],
    'NVIDIA: "Physical AI의 시대가 시작됐다"',
    ["Jensen Huang CES 2026 키노트 핵심:",
     "",
     "🧠 The Brain — 기반 지능",
     "     Cosmos Reason 2: 물리적 속성 추론",
     "",
     "🦾 The Body — 물리적 구현체",
     "     Cosmos Predict 2.5: 행동 결과 예측",
     "",
     "🌍 The World — 실제 환경",
     "     Sim-to-Real: 가상→현실 학습 전이"],
    "AI가 몸을 얻고, 물리 세계에서 감지·판단·행동하는 시대")

# --- 6. Vera Rubin ---
headline_slide("Vera Rubin — 차세대 AI 슈퍼칩",
    "6-Chip Co-design · 추론 비용 1/10 절감 · Alpamayo 오픈소스\n자율주행차, 로봇, 산업 자동화를 하나의 허브로 연결")

# --- 7. PART 1 divider ---
section_div("PART 1", "Physical AI — 지능이 몸을 얻다")

# --- 8. 로봇 혁명 ---
photo_left(P["booth"],
    "로봇 혁명 — 50+ 휴머노이드",
    ["CES 2026 웨스트 홀 = 로봇 박람회",
     "",
     "현대 × Boston Dynamics Atlas",
     "  → 2028년 제조 현장 투입 목표",
     "  → SDF(Software Defined Factory)",
     "",
     "Unitree G1 — $13,500",
     "  → 35kg, 3D LiDAR+카메라",
     "  → 휴머노이드 대중화의 시작",
     "",
     "Sanctuary AI · LG CLOiD · Doosan"])

# --- 9. 산업 Physical AI ---
photo_right(P["display"],
    "산업 Physical AI",
    ["Caterpillar — 자율 광산 운반 24시간",
     "John Deere — AI 잡초 감지·제거",
     "Kubota — 소규모 농가 AI 민주화",
     "",
     "Hesai — LiDAR 누적 100만 개 돌파",
     "  → ADAS + Robotics 이중 트랙",
     "",
     "Doosan Robotics — AI 협동 로봇",
     "Bucket Robotics — 자율 굴삭"],
    "제조, 물류, 농업, 가정까지 — AI가 물리 세계에서 작동하는 시대")

# --- 10. PART 2 divider ---
section_div("PART 2", "Software Defined Vehicle\n자동차의 OS 전환")

# --- 11. SDV 시장 ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
txt(slide, 0.8, 0.4, 11, 0.8, "SDV 시장 — iPhone/Android 모멘트", 28, TEXT, True)
stat_card(slide, 0.8, 1.5, 3.5, 2.5, "$2,135억", "2024 SDV 시장")
txt(slide, 4.6, 2.2, 0.8, 1, "→", 36, AI, True, PP_ALIGN.CENTER)
stat_card(slide, 5.5, 1.5, 3.5, 2.5, "$1조 2,400억", "2030 SDV 시장 전망")
stat_card(slide, 9.5, 1.5, 3.5, 2.5, "3.4%→90%", "SDV 비중\n(2024→2029)")
txt(slide, 0.8, 4.5, 11, 1, "자동차 산업은 지금 스마트폰이 겪었던 OS 전환을 통과하고 있다", 22, MUTED, False, PP_ALIGN.CENTER)
txt(slide, 0.8, 5.3, 11, 1, "피처폰 → iOS/Android  =  기계 중심 → 소프트웨어 중심", 20, AI, True, PP_ALIGN.CENTER)

# --- 12. 5계층 아키텍처 ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
txt(slide, 0.8, 0.3, 11, 0.8, "SDV 5계층 아키텍처", 28, TEXT, True)
layers = [
    ("L5", "Physical Embodiment", "HMG · Caterpillar · Kia PBV", 0.95),
    ("L4", "Autonomous Agent", "Waymo · Zoox · Cerence AI", 0.80),
    ("L3", "World Model & Perception", "HERE · Hesai · NVIDIA Cosmos", 0.65),
    ("L2", "Digital Nervous System", "QNX · Elektrobit · dSPACE", 0.50),
    ("L1", "Silicon Foundation", "Qualcomm · NVIDIA · Synopsys", 0.35),
]
y = 1.3
for label, name, companies, opacity in layers:
    c = RGBColor(int(108*opacity), int(99*opacity), int(255*opacity))
    box(slide, 0.8, y, 11.7, 1.0, SURFACE)
    txt(slide, 1.0, y+0.15, 1.2, 0.6, label, 22, AI, True)
    txt(slide, 2.5, y+0.15, 4, 0.6, name, 20, TEXT, True)
    txt(slide, 7.5, y+0.2, 5, 0.5, companies, 16, MUTED)
    y += 1.1
txt(slide, 0.8, 6.8, 11, 0.5, "승자 = 전체 스택을 통제하고 최적화하는 기업", 16, AI, True, PP_ALIGN.CENTER)

# --- 13. Waymo ---
photo_left(P["waymo1"],
    "Waymo — 풀스택 플랫폼 사업자",
    ["주당 25만 건+ 라이드",
     "5개 도시 서비스",
     "100억 마일+ 자율주행 데이터",
     "",
     "OS · 센서 · 매핑 · 시뮬레이션 · 운영",
     "→ 모두 통제하는 풀스택 전략",
     "",
     "Tesla와 함께 SDV 경쟁의 기준점"])

# --- 14. OEM 전략 1 ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
txt(slide, 0.8, 0.3, 11, 0.8, "완성차 OEM의 SDV 전략", 28, TEXT, True)
# 5 cards
cards = [
    ("Kia PBV", "PV5(2026) · PV7(2027)\nEasy Swap 모듈 전환\n차량 = 플랫폼"),
    ("BMW iX3", "Alexa + LLM 통합\n파노라믹 디스플레이\n자연어 차량 제어"),
    ("AFEELA", "$89,900 · Cloud AI\nPlayStation 통합\n움직이는 엔터테인먼트"),
    ("VW IDA", "ChatGPT 차량 통합\n자연어 질문 즉답\nLLM = 기본 인터페이스"),
    ("Geely/Zeekr", "GEEA 3.0 아키텍처\nXingrui LLM\n차량·모바일·홈 통합"),
]
x = 0.5
for title, desc in cards:
    box(slide, x, 1.3, 2.4, 4.5, SURFACE)
    txt(slide, x+0.15, 1.5, 2.1, 0.5, title, 18, AI, True, PP_ALIGN.CENTER)
    multi_txt(slide, x+0.15, 2.2, 2.1, 3,
        [(line, 14, MUTED, False) for line in desc.split('\n')])
    x += 2.5

# --- 15. 플랫폼 전쟁 ---
photo_right(P["floor2"],
    "플랫폼 전쟁의 주역들",
    ["Qualcomm — Snapdragon Digital Chassis",
     "  100개+ OEM · $220억 목표(2029)",
     "",
     "BlackBerry QNX — QNX Everywhere",
     "  SDP 8.0 무료 제공 · 7,500+ 고객사",
     "  ISO 26262 ASIL-D 인증",
     "",
     "Elektrobit",
     "  AMD + Android Auto + Gemini AI",
     "  + Unreal Engine + HERE 통합"])

# --- 16. PART 3 divider ---
section_div("PART 3", "현장 체험\nZoox · Waymo · Vegas Loop")

# --- 17. Zoox 외관 ---
photo_slide(P["zoox_ext"], "Zoox — 핸들 없는 차에 타다",
    "Amazon 산하 · 전용 설계 로보택시 · 핸들도 페달도 없다")

# --- 18. Zoox 내부 ---
photo_left(P["zoox_in1"],
    "Zoox 탑승 체험",
    ["양방향 좌석 · 360도 LiDAR+카메라",
     "앱으로 호출 → 자동 도착 → 슬라이딩 도어",
     "",
     "100만 마일+ 자율주행",
     "30만+ 탑승객",
     "",
     "5분 지나니 운전자가 없다는 사실을 잊었다",
     "",
     "한계: 운행 구간 제한적",
     "핵심: '완벽한 기술'이 아닌 '충분히 안전한 서비스'"])

# --- 19. Zoox 인사이트 ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
txt(slide, 0.8, 0.3, 11, 0.8, "Zoox 탑승 인사이트", 28, TEXT, True)
insights = [
    ("🧠", "Perception → Reasoning", "인식에서 추론으로 — 모호한 상황에서 판단하는 능력이 L4의 핵심"),
    ("🏙️", "도시 인프라 공진화", "기술만의 문제가 아닌 도시 설계·규제·보험 생태계 전체의 공진화"),
    ("💰", "비즈니스 모델 재편", "운전자 인건비 60~70% → 자율주행으로 비용 구조 근본 변화"),
]
y = 1.3
for icon, title, desc in insights:
    box(slide, 0.8, y, 11.7, 1.6, SURFACE)
    txt(slide, 1.2, y+0.2, 0.8, 0.8, icon, 32, TEXT, False, PP_ALIGN.CENTER)
    txt(slide, 2.2, y+0.2, 4, 0.6, title, 22, TEXT, True)
    txt(slide, 2.2, y+0.8, 9.5, 0.6, desc, 16, MUTED)
    y += 1.8

# --- 20. Zoox vs Waymo ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
txt(slide, 0.8, 0.3, 11, 0.8, "Zoox vs Waymo: 두 가지 전략", 28, TEXT, True)
ph1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(6), Inches(3))
ph1.fill.solid(); ph1.fill.fore_color.rgb = RGBColor(0x22,0x22,0x30); ph1.line.color.rgb = RGBColor(0x3A,0x3A,0x4A); ph1.line.width = Pt(1)
txt(slide, 0.7, 2.2, 5.6, 1, "📷 Zoox 로보택시 사진\n(외관 또는 측면)", 12, MUTED, False, PP_ALIGN.CENTER)
ph2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.3), Inches(6), Inches(3))
ph2.fill.solid(); ph2.fill.fore_color.rgb = RGBColor(0x22,0x22,0x30); ph2.line.color.rgb = RGBColor(0x3A,0x3A,0x4A); ph2.line.width = Pt(1)
txt(slide, 7.0, 2.2, 5.6, 1, "📷 Waymo 로보택시 사진\n(센서 어레이 또는 주행 모습)", 12, MUTED, False, PP_ALIGN.CENTER)
box(slide, 0.5, 4.5, 6, 2.5, SURFACE)
txt(slide, 0.7, 4.6, 5.6, 0.5, "Zoox — 전용 설계", 20, AI, True)
txt(slide, 0.7, 5.2, 5.6, 1.5, "핸들 없는 완전 새 설계\n탑승 경험 차별화\n스케일링은 과제", 15, MUTED)
box(slide, 6.8, 4.5, 6, 2.5, SURFACE)
txt(slide, 7.0, 4.6, 5.6, 0.5, "Waymo — 기존 차량 활용", 20, AI, True)
txt(slide, 7.0, 5.2, 5.6, 1.5, "OEM 양산 라인 활용\n심리적 진입 장벽 낮음\n주당 25만 건+ 라이드", 15, MUTED)

# --- 21. Vegas Loop ---
photo_left(P["vegas"],
    "Vegas Loop — 지하의 자율주행",
    ["The Boring Company 지하 터널 루프",
     "Tesla 차량으로 주요 지점 연결",
     "",
     "CES 기간 8만 명+ 이용",
     "10마일+ 운영 → 68마일 확장 계획",
     "리조트 $4.25 · 공항 $12",
     "",
     "지상 로보택시 + 지하 터널",
     "→ AI 통합 모빌리티 네트워크의 프로토타입"])

# --- 22. PART 4 divider ---
section_div("PART 4", "모빌리티 분야에서 주목한 기업 50")


# --- 50 Companies Infographic ---
def companies_infographic():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    txt(slide, 0.4, 0.15, 12, 0.6, "모빌리티 분야에서 주목한 기업 50", 24, TEXT, True, PP_ALIGN.CENTER)
    
    # Theme colors
    TC = {
        "ad": RGBColor(0x6C, 0x63, 0xFF),      # purple - autonomous driving
        "sdv": RGBColor(0x38, 0x9C, 0xFF),      # blue - SDV
        "os": RGBColor(0x10, 0xB9, 0x81),       # green - OS
        "ux": RGBColor(0xF5, 0x9E, 0x0B),       # amber - UX
        "robot": RGBColor(0xEF, 0x44, 0x44),     # red - robotics
        "sim": RGBColor(0xA7, 0x8B, 0xFA),      # light purple - simulation
        "chip": RGBColor(0x60, 0xA5, 0xFA),     # light blue - chips
        "sensor": RGBColor(0x34, 0xD3, 0x99),   # light green - sensors
        "industry": RGBColor(0xFB, 0xBF, 0x24), # yellow - industry
        "beyond": RGBColor(0xF8, 0x71, 0x71),   # light red - beyond
    }
    
    themes = [
        ("🚗 Autonomous Driving", "ad", ["Waymo", "Zoox", "Mobileye", "Imagry", "Tensor Auto"]),
        ("⚡ EV & SDV", "sdv", ["HMG", "Geely", "ZF", "KPIT", "HERE", "HL Mando", "Mobis", "WIA"]),
        ("🖥️ OS & Middleware", "os", ["QNX", "Elektrobit", "Green Hills", "Sonatus", "AWS"]),
        ("🎯 UX", "ux", ["Cerence", "SoundHound", "Smart Eye", "emotion3D"]),
        ("🤖 Robotics", "robot", ["Doosan", "Bucket"]),
        ("🧪 Simulation", "sim", ["Synopsys", "dSPACE", "Cognata", "Foretellix"]),
        ("🔲 AI Chips", "chip", ["NVIDIA", "Qualcomm", "Hailo"]),
        ("📡 Sensors", "sensor", ["AEye", "Hesai", "RoboSense", "LG Innotek", "BOE", "Aeva", "Arbe", "Innoviz", "Lidwave"]),
        ("🏗️ Industry AI", "industry", ["Deere", "CAT", "Kubota", "Brunswick", "Oshkosh"]),
        ("🔋 Beyond", "beyond", ["Aptiv", "Valeo", "Electra", "Autofleet", "Donut Lab"]),
    ]
    
    # Layout: 2 columns of theme groups
    col_x = [0.3, 6.8]
    col_themes = [themes[:5], themes[5:]]
    
    for col_idx, col_data in enumerate(col_themes):
        x_base = col_x[col_idx]
        y = 0.85
        for theme_label, theme_key, companies in col_data:
            color = TC[theme_key]
            # Theme label
            txt(slide, x_base, y, 6, 0.35, theme_label, 11, color, True)
            y += 0.32
            # Company bubbles
            bx = x_base
            for company in companies:
                w = max(0.7, len(company) * 0.1 + 0.3)
                if bx + w > x_base + 6.2:
                    y += 0.35
                    bx = x_base
                s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(bx), Inches(y), Inches(w), Inches(0.3))
                s.fill.solid()
                s.fill.fore_color.rgb = RGBColor(0x15, 0x15, 0x20)
                s.line.color.rgb = color
                s.line.width = Pt(0.75)
                s.shadow.inherit = False
                tf = s.text_frame
                tf.word_wrap = False
                p = tf.paragraphs[0]
                p.text = company
                p.font.size = Pt(9)
                p.font.color.rgb = color
                p.font.bold = True
                p.font.name = 'Arial'
                p.alignment = PP_ALIGN.CENTER
                bx += w + 0.08
            y += 0.55
    
    # Legend at bottom
    txt(slide, 0.5, 7.0, 12, 0.3, "■ Key Theme    ■ Sub Theme", 10, MUTED, False, PP_ALIGN.CENTER)

companies_infographic()

# --- 23-26. 기업 50 (4 slides, compact) ---
photo_right(P["booth"],
    "Key Themes — 자율주행 & SDV",
    ["🚗 Waymo · Zoox · Mobileye · Imagry · Tensor Auto",
     "",
     "⚡ HMG · Geely · ZF · KPIT · HERE",
     "   HL Mando · Mobis · WIA",
     "",
     "🖥️ QNX · Elektrobit · Green Hills · Sonatus · AWS",
     "",
     "🎯 Cerence · SoundHound · Smart Eye · emotion3D",
     "",
     "🤖 Doosan Robotics · Bucket Robotics"])

photo_left(P["street1"],
    "Sub Themes — 시뮬레이션 · 칩 · 센서",
    ["🧪 Synopsys · dSPACE · Cognata · Foretellix",
     "",
     "🔲 NVIDIA · Qualcomm · Hailo",
     "",
     "📡 AEye · Hesai · RoboSense · LG Innotek",
     "   BOE · Aeva · Arbe · Innoviz · Lidwave"])

photo_right(P["street2"],
    "Sub Themes — 산업 AI · Beyond Mobility",
    ["🏗️ John Deere · Caterpillar · Kubota",
     "   Brunswick · Oshkosh",
     "",
     "🔋 Aptiv · Valeo · Electra Vehicles",
     "   Autofleet · Donut Lab"])

# --- 27. 종합 인사이트 divider ---
section_div("종합 인사이트", "CES 2026이 남긴 다섯 가지 메시지")

# --- 28-32. 5 Key Messages ---
key_msg_slide(1, "AI는 이제 행동한다",
    "움직이고, 감지하고, 행동한다.\n자율주행차, 로봇, 산업 기계가 공통의\n'구현 지능 스택'을 공유하기 시작했다.")

key_msg_slide(2, "부품에서 플랫폼으로",
    "개별 센서나 칩이 아니라\n엔드-투-엔드 플랫폼을 누가 장악하느냐가\n승부를 가른다.")

key_msg_slide(3, "SDV는 OS 전쟁이다",
    "미들웨어, 소프트웨어 아키텍처,\nOTA 업데이트 역량이\nOEM의 미래 경쟁력을 결정한다.")

key_msg_slide(4, "인간 중심 UX가 차별화",
    "안전, 신뢰, 설명 가능성, 직관적 사용성이\n소프트웨어 진화의 완성점이다.\n'차가 나를 이해하는' 경험이 브랜드를 결정한다.")

key_msg_slide(5, "Physical AI = 문명 인프라",
    "제조, 에너지, 헬스케어, 도시, 농업까지 —\nAI는 모든 산업 변혁의 핵심 엔진이다.\n모빌리티를 넘어 문명의 인프라가 된다.")

# --- 33. 결론 ---
photo_slide(P["night"],
    "승자는 플랫폼을 통제하고,\n소프트웨어를 장악하며,\n인간 경험을 설계하는 자다.", "", 36)

# --- 34. 비즈니스 임플리케이션 ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
txt(slide, 0.8, 0.3, 11, 0.8, "비즈니스 임플리케이션", 28, TEXT, True)
actions = [
    ("01", "Physical AI 투자 가속화", "로봇·자율주행·산업 자동화 예산 확대"),
    ("02", "SDV 전환 대비", "소프트웨어 역량 내재화 또는 파트너십 확보"),
    ("03", "플랫폼 전략 수립", "단일 기술이 아닌 스택 통합 관점"),
    ("04", "인간 중심 설계", "UX/AI 인터페이스 투자 우선순위 상향"),
    ("05", "생태계 포지셔닝", "5계층 아키텍처에서 자사의 위치 정의"),
]
y = 1.3
for num, title, desc in actions:
    box(slide, 0.8, y, 11.7, 1.0, SURFACE)
    txt(slide, 1.2, y+0.15, 0.8, 0.6, num, 22, AI, True)
    txt(slide, 2.2, y+0.1, 3.5, 0.6, title, 20, TEXT, True)
    txt(slide, 6.0, y+0.15, 6, 0.6, desc, 16, MUTED)
    y += 1.1

# --- 35. Q&A ---
photo_slide(P["pano"], "Q & A", "감사합니다\n\n우지환 · AWS Senior Specialist Partner SA — AI/ML", 52)

# --- 36. Appendix ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
txt(slide, 0.8, 0.4, 11, 0.8, "참고 자료", 28, TEXT, True)
sources = [
    "CES 2026 공식: ces.tech",
    "NVIDIA CES 키노트: blogs.nvidia.com",
    "S&P Global Automotive: spglobal.com",
    "Waymo: waymo.com · Zoox: zoox.com",
    "Vegas Loop: vegasloop.com",
    "Kia PBV: kia.com · AFEELA: shm-afeela.com",
    "",
    "본 보고서는 CES 2026 현장 참관 및 수집 자료를 바탕으로 작성되었습니다.",
]
y = 1.5
for s in sources:
    if s:
        txt(slide, 1.0, y, 11, 0.4, s, 16, MUTED)
    y += 0.4

# ============================================================
# SAVE
# ============================================================
output = os.path.join(BASE, 'CES_2026_Report.pptx')
prs.save(output)
print(f"✅ PPT saved: {output}")
print(f"   Total slides: {len(prs.slides)}")
